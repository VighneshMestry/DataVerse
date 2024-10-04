# docx==0.8.11
# python-pptx==0.6.18
# openpyxl==3.0.9
# PyMuPDF==1.19.6
# requests==2.26.0
# pytesseract==0.3.8
# Pillow==8.4.0
# opencv-python==4.5.3.20210920
# numpy==1.21.2


from io import BytesIO
import sys
from docx import Document
from pptx import Presentation
import openpyxl # For excel
import fitz
import requests
import pytesseract
from PIL import Image

def read_pdf_from_url(url):
    # Fetch the content of the PDF from the URL
    response = requests.get(url)
    # pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
    
    # Check if the request was successful
    if response.status_code == 200:
        # Initialize an empty string to store the text
        text = ''
        # Open the PDF file using PyMuPDF
        pdf_document = fitz.open(stream=response.content, filetype="pdf")
        
        # Iterate through each page and extract text
        for page_num in range(pdf_document.page_count):
            page = pdf_document.load_page(page_num)
            
            images = page.get_images(full=True)
            if(images):
                # Iterate through each image on the page
                for img_info in images:
                    img_index = img_info[0]  # Extracting the image index
                    img = page.get_pixmap()
                    img_pil = Image.frombytes("RGB", (img.width, img.height), img.samples)
                    myconfig = r"--psm 6 --oem 3"
                    text += pytesseract.image_to_string(img_pil, lang='eng', config=myconfig)
            else:
                text += page.get_text().encode("utf-8").decode("utf-8") 
        text = text.encode('utf-8', 'ignore').decode('utf-8') 
        cleaned_str = remove_non_ascii_chars(text)
        return cleaned_str.encode('utf-8')
    else:
        # print(f"Failed to fetch PDF from URL. Status code: {response.status_code}")
        return "None"
    
def read_online_word_document(url):
    # Send a GET request to fetch the Word document
    response = requests.get(url)
    
    # Check if the request was successful (status code 200)
    if response.status_code == 200:
        # Read the content of the Word document
        docx_content = response.content
        
        # Create a BytesIO object to work with in-memory binary data
        docx_file = BytesIO(docx_content)
        
        # Load the Word document using python-docx
        doc = Document(docx_file)
        
        # Extract text from the Word document
        text = ""
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        cleaned_str = remove_non_ascii_chars(text)
        return cleaned_str.encode('utf-8')
    else:
        # print(f"Failed to fetch Word document. Status code: {response.status_code}")
        return None
    
def read_online_ppt_document(url):
    # Send a GET request to fetch the Word document
    response = requests.get(url)
    
    # Check if the request was successful (status code 200)
    if response.status_code == 200:
        # Read the content of the Word document
        ppt_content = response.content
        
        # Create a BytesIO object to work with in-memory binary data
        ppt_file = BytesIO(ppt_content)
        
        # Load the Word document using python-docx
        presentation = Presentation(ppt_file)
        
        # Extract text from the Word document
        text = ""
        # for paragraph in doc.paragraphs:
        #     text += paragraph.text + "\n"
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text+=shape.text
        
        cleaned_str = remove_non_ascii_chars(text)
        return cleaned_str.encode('utf-8')
    else:
        # print(f"Failed to fetch Word document. Status code: {response.status_code}")
        return None
    
def read_online_xlsx_document(url):
    # Send a GET request to fetch the Word document
    response = requests.get(url)
    
    # Check if the request was successful (status code 200)
    if response.status_code == 200:
        # Read the content of the Word document
        xlsx_content = response.content
        
        # Create a BytesIO object to work with in-memory binary data
        xlsx_file = BytesIO(xlsx_content)
        
        # Load the Word document using python-docx
        wb = openpyxl.load_workbook(xlsx_file)
        
        # Extract text from the Word document
        text = ""
        # for paragraph in doc.paragraphs:
        #     text += paragraph.text + "\n"
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            for row in sheet.iter_rows(values_only=True):
                text += ' '.join(str(cell) for cell in row)
        
        cleaned_str = remove_non_ascii_chars(text)
        return cleaned_str.encode('utf-8')
    else:
        # print(f"Failed to fetch Word document. Status code: {response.status_code}")
        return None


def image_to_text(url):
    try:
        # Download the image from the URL
        # pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'
        response = requests.get(url)
        response.raise_for_status()  # Raise an exception if download fails

        # Open the image from binary data
        with Image.open(BytesIO(response.content)) as img:
            # Use pytesseract to do OCR(Optical Character Recognition) on the image
            text = pytesseract.image_to_string(img)
            return text.strip()  # Remove leading/trailing whitespace
    except Exception as e:
        print(f"Error reading text from image: {e}")
        return None

def remove_non_ascii_chars(input_str):
    """
    Remove non-ASCII characters from the input string.
    """
    return ''.join(char for char in input_str if ord(char) < 128)


if __name__ == "__main__":
    document_path = sys.argv[1]
    # print("Document path is " + document_path)  # Replace with the path to your document
    
    # # Split the filename by '.' and get the last part as the file extension
    extracted_text = ""

    filename = document_path.split('/')[-1]
    
    # Split the filename by '.' and get the last part as the file extension
    file_extension = filename.split('.')[-1]
    file_extension = file_extension[0:4]
    
    # file_extension =  file_extension.lower()
    # print("File extension is " + file_extension)
    if file_extension == 'pdf?':
        extracted_text = read_pdf_from_url(document_path)
    elif file_extension == 'docx':
        extracted_text = read_online_word_document(document_path)
    elif file_extension == 'pptx':
        extracted_text = read_online_ppt_document(document_path)
    elif file_extension == 'xlsx':
        extracted_text = read_online_xlsx_document(document_path)
    elif file_extension == 'jpg?':
        # print("The URL points to an Excel file.")
        extracted_text = image_to_text(document_path)
    elif file_extension == 'png?':
        # print("The URL points to an Excel file.")
        extracted_text = image_to_text(document_path)
    elif file_extension == 'docs':
        # print("The URL points to an Excel file.")
        extracted_text = image_to_text(document_path)
    else:
        # print("The file type is unknown.")
        error = ''

    print(extracted_text)